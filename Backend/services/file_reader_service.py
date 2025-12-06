import os
import json
import yaml
import defusedxml.ElementTree as ET
import io
import csv
import magic

from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from odf.opendocument import load as odf_load
from odf.text import P
from odf import teletype
from PyPDF2 import PdfReader

# Limite de caracteres para extração de texto (aprox 10MB de texto puro)
MAX_TEXT_CHARS = 10 * 1024 * 1024

"""
Formatos suportados:
- TXT, MD, CSV, PDF, RTF, DOC, DOCX, ODT, PPTX, XLSX, HTML, XML, JSON, YAML, YML
"""

def validate_file_type(file_stream, expected_ext):
    """
    Valida o tipo real do arquivo usando Magic Numbers.
    Retorna True se o arquivo for seguro/correspondente.
    """
    # Lê os primeiros 2KB para identificar o cabeçalho
    header = file_stream.read(2048)
    file_stream.seek(0) # Reseta o ponteiro para o início
    
    mime = magic.Magic(mime=True)
    detected_mime = mime.from_buffer(header)
    
    # Mapeamento simplificado de extensões permitidas para Mime Types
    # Adicione mais conforme necessário
    valid_mimes = {
        'pdf': ['application/pdf'],
        'docx': ['application/vnd.openxmlformats-officedocument.wordprocessingml.document'],
        'txt': ['text/plain'],
        'csv': ['text/plain', 'text/csv'],
        'json': ['application/json', 'text/plain'],
        'xml': ['text/xml', 'application/xml'],
        'yaml': ['application/x-yaml', 'text/plain'],
        'yml': ['application/x-yaml', 'text/plain'],
        'rtf': ['application/rtf', 'text/rtf'],
        'odt': ['application/vnd.oasis.opendocument.text'],
        'pptx': ['application/vnd.openxmlformats-officedocument.presentationml.presentation'],
        'xlsx': ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'],
        'md': ['text/markdown', 'text/plain'],
        'html': ['text/html'],
    }
    
    if expected_ext in valid_mimes:
        # Se o mime detectado estiver na lista de permitidos ou for genérico (octet-stream as vezes acontece com arquivos novos)
        if detected_mime in valid_mimes[expected_ext]:
            return True
        # Fallback para text/plain que muitas vezes é detectado incorretamente
        if 'text' in detected_mime and 'text' in valid_mimes[expected_ext][0]:
            return True
            
    # Se não conseguimos validar estritamente, mas não é um executável perigoso
    blocked_mimes = ['application/x-dosexec', 'application/x-executable']
    if detected_mime in blocked_mimes:
        raise ValueError(f"Arquivo malicioso detectado. Tipo real: {detected_mime}")
        
    return True

def extract_text_from_file(file_stream, filename: str) -> str:
    """Extrai texto de praticamente qualquer formato relevante para contratos."""

    ext = os.path.splitext(filename)[1].lower().replace(".", "")
    
    # Validação de segurança antes de processar
    try:
        validate_file_type(file_stream, ext)
    except Exception as e:
        return f"Erro de segurança: {str(e)}"

    if ext == "txt" or ext == "md":
        # CORREÇÃO B: Limitar leitura para evitar estouro de memória
        content = file_stream.read(MAX_TEXT_CHARS).decode("utf-8", errors="ignore")
        return content

    if ext == "csv":
        return _extract_csv(file_stream)

    if ext == "pdf":
        return _extract_pdf(file_stream)

    if ext == "rtf":
        return _extract_rtf(file_stream)

    if ext in ["doc", "docx"]:
        return _extract_docx(file_stream)

    if ext == "odt":
        return _extract_odt(file_stream)

    if ext == "pptx":
        return _extract_pptx(file_stream)

    if ext == "xlsx":
        return _extract_xlsx(file_stream)

    if ext == "html":
        return _extract_html(file_stream)

    if ext == "xml":
        return _extract_xml(file_stream)

    if ext == "json":
        return json.dumps(json.load(file_stream), ensure_ascii=False, indent=2)

    if ext in ["yaml", "yml"]:
        return yaml.safe_dump(yaml.safe_load(file_stream), allow_unicode=True)

    raise ValueError(f"Formato não suportado: {ext}")


def _extract_pdf(fs):
    reader = PdfReader(fs)
    text = []
    for page in reader.pages:
        t = page.extract_text()
        if t:
            text.append(t)
    return "\n".join(text)


def _extract_csv(fs):
    fs.seek(0)
    decoded = io.StringIO(fs.read().decode("utf-8", errors="ignore"))
    reader = csv.reader(decoded)
    return "\n".join([", ".join(row) for row in reader])


def _extract_rtf(fs):
    try:
        import striprtf
    except ImportError:
        return "Erro: A biblioteca 'striprtf' é necessária para extrair RTF."

    return striprtf.rtf_to_text(fs.read().decode("utf-8", errors="ignore"))


def _extract_docx(fs):
    doc = Document(fs)
    return "\n".join([p.text for p in doc.paragraphs])


def _extract_odt(fs):
    odt = odf_load(fs)
    paragraphs = odt.getElementsByType(P)
    return "\n".join([teletype.extractText(p) for p in paragraphs])


def _extract_pptx(fs):
    prs = Presentation(fs)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def _extract_xlsx(fs):
    wb = load_workbook(fs, data_only=True)
    textos = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_txt = [str(cell) for cell in row if cell is not None]
            if row_txt:
                textos.append(" | ".join(row_txt))
    return "\n".join(textos)


def _extract_html(fs):
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(fs.read(), "html.parser")
    return soup.get_text(separator="\n")


def _extract_xml(fs):
    tree = ET.parse(fs)
    root = tree.getroot()
    return _xml_recursive_text(root)


def _xml_recursive_text(node):
    text = (node.text or "").strip()
    for child in node:
        text += "\n" + _xml_recursive_text(child)
    return text
